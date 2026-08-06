"""Pure-Python fallbacks for DOCX / PPTX / XLSX / CSV.

These kick in when both Docling and MarkItDown are unavailable (or fail),
or when the user explicitly disables ML-based extractors. They produce
clean Markdown using only wheel-only Python libraries — so they work on
all OSes, including stripped-down Windows installs.
"""

from __future__ import annotations

import csv
import time
from io import StringIO
from pathlib import Path

from convert_anything_md.extractors.base import (
    ExtractionResult,
    ExtractorError,
    ExtractorUnavailable,
    word_count,
)


class DocxExtractor:
    """DOCX → Markdown via python-docx. Preserves paragraphs + tables + headings."""

    name = "python-docx"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            from docx import Document  # python-docx
            from docx.document import Document as _DocumentType  # noqa: F401
        except ImportError as exc:
            raise ExtractorUnavailable("python-docx is not installed") from exc

        start = time.perf_counter()
        try:
            doc = Document(str(path))
        except Exception as exc:  # noqa: BLE001
            raise ExtractorError(f"python-docx failed on {path.name}: {exc}") from exc

        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading"):
                # "Heading 1" → "# …", "Heading 2" → "## …", etc.
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 2
                level = max(1, min(6, level))
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)

        for table in doc.tables:
            rows = [
                [cell.text.strip().replace("\n", " ") for cell in row.cells]
                for row in table.rows
            ]
            if rows:
                parts.append(_render_markdown_table(rows))

        markdown = "\n\n".join(parts).strip() + "\n"
        duration_ms = int((time.perf_counter() - start) * 1000)

        return ExtractionResult(
            markdown=markdown,
            engine=self.name,
            word_count=word_count(markdown),
            duration_ms=duration_ms,
            fallback_chain=[self.name],
        )


class PptxExtractor:
    """PPTX → Markdown via python-pptx. One H2 section per slide."""

    name = "python-pptx"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorUnavailable("python-pptx is not installed") from exc

        start = time.perf_counter()
        try:
            prs = Presentation(str(path))
        except Exception as exc:  # noqa: BLE001
            raise ExtractorError(f"python-pptx failed on {path.name}: {exc}") from exc

        sections: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines = [f"## Slide {idx}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                text = text.strip()
                if text:
                    lines.append(text)
            # Slide notes if present.
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide is not None:
                note_text = (notes_slide.notes_text_frame.text or "").strip()
                if note_text:
                    lines.append(f"**Notes:** {note_text}")
            sections.append("\n\n".join(lines))

        markdown = "\n\n---\n\n".join(sections).strip() + "\n"
        duration_ms = int((time.perf_counter() - start) * 1000)

        return ExtractionResult(
            markdown=markdown,
            engine=self.name,
            page_count=len(prs.slides),
            word_count=word_count(markdown),
            duration_ms=duration_ms,
            fallback_chain=[self.name],
        )


class XlsxExtractor:
    """XLSX → Markdown via openpyxl. One section per sheet, table per sheet."""

    name = "openpyxl"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise ExtractorUnavailable("openpyxl is not installed") from exc

        start = time.perf_counter()
        try:
            wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        except Exception as exc:  # noqa: BLE001
            raise ExtractorError(f"openpyxl failed on {path.name}: {exc}") from exc

        sections: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([
                    "" if cell is None else str(cell).replace("\n", " ").strip()
                    for cell in row
                ])
            # Trim trailing all-empty rows.
            while rows and not any(rows[-1]):
                rows.pop()
            sections.append(f"## {sheet_name}")
            if rows:
                sections.append(_render_markdown_table(rows))
            else:
                sections.append("_(empty sheet)_")

        wb.close()
        markdown = "\n\n".join(sections).strip() + "\n"
        duration_ms = int((time.perf_counter() - start) * 1000)

        return ExtractionResult(
            markdown=markdown,
            engine=self.name,
            word_count=word_count(markdown),
            duration_ms=duration_ms,
            fallback_chain=[self.name],
        )


class CsvExtractor:
    """CSV / TSV → Markdown table. Pure stdlib, no deps."""

    name = "csv"

    def extract(self, path: Path) -> ExtractionResult:
        start = time.perf_counter()
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","

        try:
            text = path.read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            raise ExtractorError(f"cannot read {path.name}: {exc}") from exc

        # Use csv.Sniffer to pick up odd delimiters in .csv files.
        if delimiter == "," and text:
            try:
                dialect = csv.Sniffer().sniff(text[:4096])
                delimiter = dialect.delimiter
            except csv.Error:
                pass

        reader = csv.reader(StringIO(text), delimiter=delimiter)
        rows = [
            [cell.strip().replace("\n", " ") for cell in row]
            for row in reader
            if row
        ]
        markdown = (
            _render_markdown_table(rows) + "\n"
            if rows
            else f"# {path.stem}\n\n_(empty file)_\n"
        )
        duration_ms = int((time.perf_counter() - start) * 1000)

        return ExtractionResult(
            markdown=markdown,
            engine=self.name,
            word_count=word_count(markdown),
            duration_ms=duration_ms,
            fallback_chain=[self.name],
            extra={"rows": len(rows), "delimiter": delimiter},
        )


def _render_markdown_table(rows: list[list[str]]) -> str:
    """Render a 2-D list as a GitHub-flavored Markdown table."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    # Pad short rows.
    padded = [list(r) + [""] * (width - len(r)) for r in rows]

    def _escape(cell: str) -> str:
        # Escape pipes so they don't break the table.
        return cell.replace("|", "\\|")

    header = padded[0]
    body = padded[1:]

    lines = [
        "| " + " | ".join(_escape(c) for c in header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines.extend(
        "| " + " | ".join(_escape(c) for c in row) + " |" for row in body
    )
    return "\n".join(lines)

